"""ONtown 사업계획서 PPT 생성 스크립트"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# 브랜드 컬러
PRIMARY = RGBColor(0x1A, 0x56, 0xDB)      # 딥 블루
ACCENT = RGBColor(0x06, 0xB6, 0xD4)        # 시안
DARK = RGBColor(0x1E, 0x29, 0x3B)          # 슬레이트 다크
GRAY = RGBColor(0x64, 0x74, 0x8B)          # 슬레이트 그레이
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_BG = RGBColor(0xF8, 0xFA, 0xFC)


def set_slide_bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_accent_bar(slide, top=Inches(0), height=Inches(0.08)):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, top, Inches(13.333), height)
    bar.fill.solid()
    bar.fill.fore_color.rgb = PRIMARY
    bar.line.fill.background()


def add_footer(slide, text="ONtown | 하이퍼로컬 커뮤니티 플랫폼"):
    box = slide.shapes.add_textbox(Inches(0.5), Inches(7.0), Inches(12), Inches(0.4))
    tf = box.text_frame
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(9)
    p.font.color.rgb = GRAY
    p.alignment = PP_ALIGN.RIGHT


def add_title_slide(prs, title, subtitle, tagline=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK)
    add_accent_bar(slide, top=Inches(6.8), height=Inches(0.2))

    # 로고 영역
    logo = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.2), Inches(1.2), Inches(1.2))
    logo.fill.solid()
    logo.fill.fore_color.rgb = ACCENT
    logo.line.fill.background()
    lt = logo.text_frame
    lt.paragraphs[0].text = "ON"
    lt.paragraphs[0].font.size = Pt(28)
    lt.paragraphs[0].font.bold = True
    lt.paragraphs[0].font.color.rgb = WHITE
    lt.paragraphs[0].alignment = PP_ALIGN.CENTER
    lt.vertical_anchor = MSO_ANCHOR.MIDDLE

    tb = slide.shapes.add_textbox(Inches(0.8), Inches(2.8), Inches(11.5), Inches(2))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = WHITE

    p2 = tf.add_paragraph()
    p2.text = subtitle
    p2.font.size = Pt(20)
    p2.font.color.rgb = ACCENT
    p2.space_before = Pt(12)

    if tagline:
        p3 = tf.add_paragraph()
        p3.text = tagline
        p3.font.size = Pt(14)
        p3.font.color.rgb = GRAY
        p3.space_before = Pt(24)


def add_section_slide(prs, section_num, section_title):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, PRIMARY)

    num_box = slide.shapes.add_textbox(Inches(0.8), Inches(2.5), Inches(2), Inches(1))
    np = num_box.text_frame.paragraphs[0]
    np.text = f"{section_num:02d}"
    np.font.size = Pt(72)
    np.font.bold = True
    np.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(3.8), Inches(11), Inches(1.5))
    tp = title_box.text_frame.paragraphs[0]
    tp.text = section_title
    tp.font.size = Pt(36)
    tp.font.bold = True
    tp.font.color.rgb = WHITE


def add_content_slide(prs, title, bullets, subtitle=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_accent_bar(slide)

    tb = slide.shapes.add_textbox(Inches(0.7), Inches(0.35), Inches(12), Inches(0.8))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = DARK

    start_y = 1.3
    if subtitle:
        sb = slide.shapes.add_textbox(Inches(0.7), Inches(1.15), Inches(12), Inches(0.5))
        sp = sb.text_frame.paragraphs[0]
        sp.text = subtitle
        sp.font.size = Pt(14)
        sp.font.color.rgb = GRAY
        start_y = 1.65

    body = slide.shapes.add_textbox(Inches(0.7), Inches(start_y), Inches(12), Inches(5.2))
    btf = body.text_frame
    btf.word_wrap = True

    for i, item in enumerate(bullets):
        if isinstance(item, tuple):
            text, level = item
        else:
            text, level = item, 0
        para = btf.paragraphs[0] if i == 0 else btf.add_paragraph()
        para.text = text
        para.level = level
        para.font.size = Pt(16 if level == 0 else 14)
        para.font.color.rgb = DARK if level == 0 else GRAY
        para.space_after = Pt(10)
        if level == 0 and i > 0:
            para.space_before = Pt(6)

    add_footer(slide)
    return slide


def add_two_column_slide(prs, title, left_title, left_items, right_title, right_items):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_accent_bar(slide)

    tb = slide.shapes.add_textbox(Inches(0.7), Inches(0.35), Inches(12), Inches(0.8))
    tb.text_frame.paragraphs[0].text = title
    tb.text_frame.paragraphs[0].font.size = Pt(28)
    tb.text_frame.paragraphs[0].font.bold = True
    tb.text_frame.paragraphs[0].font.color.rgb = DARK

    for col, (col_title, items, x) in enumerate([
        (left_title, left_items, 0.7),
        (right_title, right_items, 6.8),
    ]):
        ct = slide.shapes.add_textbox(Inches(x), Inches(1.2), Inches(5.8), Inches(0.5))
        ct.text_frame.paragraphs[0].text = col_title
        ct.text_frame.paragraphs[0].font.size = Pt(18)
        ct.text_frame.paragraphs[0].font.bold = True
        ct.text_frame.paragraphs[0].font.color.rgb = PRIMARY

        body = slide.shapes.add_textbox(Inches(x), Inches(1.7), Inches(5.8), Inches(5))
        btf = body.text_frame
        btf.word_wrap = True
        for i, item in enumerate(items):
            para = btf.paragraphs[0] if i == 0 else btf.add_paragraph()
            para.text = f"• {item}"
            para.font.size = Pt(14)
            para.font.color.rgb = DARK
            para.space_after = Pt(8)

    add_footer(slide)


def add_roadmap_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_accent_bar(slide)

    tb = slide.shapes.add_textbox(Inches(0.7), Inches(0.35), Inches(12), Inches(0.8))
    tb.text_frame.paragraphs[0].text = "성장 전략 및 로드맵"
    tb.text_frame.paragraphs[0].font.size = Pt(28)
    tb.text_frame.paragraphs[0].font.bold = True
    tb.text_frame.paragraphs[0].font.color.rgb = DARK

    phases = [
        ("1단계", "0~12개월", "거점 확보", [
            "강남on 플래그십 런칭",
            "여수·강릉·부산on 우선 개설",
            "'출장·여행 전 필수 앱' 포지셔닝",
        ]),
        ("2단계", "12~24개월", "클랜 문화 형성", [
            "지역 대항 테트리스 리그",
            "얼짱 대회 지역 예선",
            "지역 태그 고도화 + 광역 커뮤니티",
            "지역 광고 상품 정식 출시",
        ]),
        ("3단계", "24개월~", "전국 플랫폼", [
            "전국 대회 정례화",
            "기획사·브랜드 제휴 확대",
            "지자체 B2G 사업 본격화",
            "커머스 연계 수익 다각화",
        ]),
    ]

    for i, (phase, period, name, items) in enumerate(phases):
        x = Inches(0.7 + i * 4.1)
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(1.3), Inches(3.8), Inches(5.2))
        card.fill.solid()
        card.fill.fore_color.rgb = LIGHT_BG
        card.line.color.rgb = RGBColor(0xE2, 0xE8, 0xF0)

        header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, Inches(1.3), Inches(3.8), Inches(0.9))
        header.fill.solid()
        header.fill.fore_color.rgb = PRIMARY if i == 0 else (ACCENT if i == 1 else DARK)
        header.line.fill.background()

        ht = header.text_frame
        ht.paragraphs[0].text = f"{phase}  {period}"
        ht.paragraphs[0].font.size = Pt(13)
        ht.paragraphs[0].font.bold = True
        ht.paragraphs[0].font.color.rgb = WHITE
        ht.paragraphs[0].alignment = PP_ALIGN.CENTER
        ht.vertical_anchor = MSO_ANCHOR.MIDDLE

        nt = slide.shapes.add_textbox(x + Inches(0.2), Inches(2.3), Inches(3.4), Inches(0.5))
        nt.text_frame.paragraphs[0].text = name
        nt.text_frame.paragraphs[0].font.size = Pt(16)
        nt.text_frame.paragraphs[0].font.bold = True
        nt.text_frame.paragraphs[0].font.color.rgb = DARK

        bt = slide.shapes.add_textbox(x + Inches(0.2), Inches(2.9), Inches(3.4), Inches(3.4))
        btf = bt.text_frame
        btf.word_wrap = True
        for j, item in enumerate(items):
            para = btf.paragraphs[0] if j == 0 else btf.add_paragraph()
            para.text = f"• {item}"
            para.font.size = Pt(12)
            para.font.color.rgb = GRAY
            para.space_after = Pt(6)

    add_footer(slide)


def add_funds_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_accent_bar(slide)

    tb = slide.shapes.add_textbox(Inches(0.7), Inches(0.35), Inches(12), Inches(0.8))
    tb.text_frame.paragraphs[0].text = "투자 유치 목적 (Use of Funds)"
    tb.text_frame.paragraphs[0].font.size = Pt(28)
    tb.text_frame.paragraphs[0].font.bold = True
    tb.text_frame.paragraphs[0].font.color.rgb = DARK

    funds = [
        ("40%", "플랫폼 개발", "지역 태그·클랜·대회 운영 시스템", PRIMARY),
        ("30%", "초기 지역 활성화", "거점 4곳 마케팅·시드 유저", ACCENT),
        ("20%", "첫 전국 이벤트", "게임 리그·얼짱 대회 파일럿", RGBColor(0x8B, 0x5C, 0xF6)),
        ("10%", "운영·제휴", "기획사·지자체·소상공인 네트워크", RGBColor(0xF5, 0x9E, 0x0B)),
    ]

    for i, (pct, title, desc, color) in enumerate(funds):
        y = Inches(1.3 + i * 1.35)
        bar_bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.7), y, Inches(11.5), Inches(1.1))
        bar_bg.fill.solid()
        bar_bg.fill.fore_color.rgb = LIGHT_BG
        bar_bg.line.fill.background()

        pct_w = float(pct.replace("%", "")) / 100 * 11.5
        bar_fill = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.7), y, Inches(pct_w), Inches(1.1))
        bar_fill.fill.solid()
        bar_fill.fill.fore_color.rgb = color
        bar_fill.line.fill.background()

        label = slide.shapes.add_textbox(Inches(1.0), y + Inches(0.15), Inches(10), Inches(0.8))
        lp = label.text_frame.paragraphs[0]
        lp.text = f"{pct}  {title}  —  {desc}"
        lp.font.size = Pt(15)
        lp.font.bold = True
        lp.font.color.rgb = WHITE if pct_w > 3 else DARK

    add_footer(slide)


def add_vision_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK)
    add_accent_bar(slide, top=Inches(6.8), height=Inches(0.2))

    tb = slide.shapes.add_textbox(Inches(0.8), Inches(1.0), Inches(11.5), Inches(1))
    tb.text_frame.paragraphs[0].text = "비전 (Vision)"
    tb.text_frame.paragraphs[0].font.size = Pt(32)
    tb.text_frame.paragraphs[0].font.bold = True
    tb.text_frame.paragraphs[0].font.color.rgb = ACCENT

    body = slide.shapes.add_textbox(Inches(0.8), Inches(2.2), Inches(11.5), Inches(3))
    btf = body.text_frame
    btf.word_wrap = True
    p = btf.paragraphs[0]
    p.text = (
        "대한민국의 모든 지역이 각자의 정체성을 가진 클랜으로 살아 숨 쉬고, "
        "서로 경쟁하고 교류하며, 그 과정에서 지역 상권이 살아나고 "
        "지역의 인재가 전국 무대로 올라서는 세계"
    )
    p.font.size = Pt(20)
    p.font.color.rgb = WHITE
    p.line_spacing = 1.5

    quote = slide.shapes.add_textbox(Inches(0.8), Inches(5.0), Inches(11.5), Inches(1.2))
    qp = quote.text_frame.paragraphs[0]
    qp.text = '"당신이 어디에 있든, 그 지역은 ON 되어 있다."'
    qp.font.size = Pt(24)
    qp.font.bold = True
    qp.font.color.rgb = ACCENT
    qp.alignment = PP_ALIGN.CENTER

    brand = slide.shapes.add_textbox(Inches(0.8), Inches(6.0), Inches(11.5), Inches(0.5))
    bp = brand.text_frame.paragraphs[0]
    bp.text = "— ONtown"
    bp.font.size = Pt(18)
    bp.font.color.rgb = GRAY
    bp.alignment = PP_ALIGN.CENTER


def build_presentation():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # 1. 표지
    add_title_slide(
        prs,
        "ONtown (온타운)",
        "하이퍼로컬 커뮤니티 플랫폼",
        "대한민국의 모든 지역을 연결하다",
    )

    # 2. 개요
    add_section_slide(prs, 1, "개요")
    add_content_slide(prs, "Executive Summary", [
        "ONtown은 '강남on', '여수on', '부산on' 등 지역 단위 커뮤니티를 전국으로 확장하는 하이퍼로컬 플랫폼",
        "각 지역의 'on'은 단순 게시판이 아닌, 지역 정체성을 가진 하나의 '클랜(Clan)'",
        "지역 간 교류·경쟁·연합을 통해 전국 단위 네트워크로 성장",
        ("핵심 차별점: 지역과 지역을 '연결'하는 것 자체가 서비스의 핵심 가치", 0),
        ("강남 유저의 여수on 댓글 → '강남' 태그 → 외지인 수요 신호 → 상권 마케팅", 1),
        ("지역 대항 e스포츠 + 전국 얼짱 대회 → 클랜 소속감 극대화", 1),
    ])

    # 3. 문제 정의
    add_section_slide(prs, 2, "문제 정의")
    add_content_slide(prs, "Problem — 왜 지금인가?", [
        ("① 지역 커뮤니티는 고립되어 있다", 0),
        ("기존 서비스는 해당 지역 거주자만 대상 → 여행자·출장자·이주 예정자 접근 불가", 1),
        ("포털·블로그는 광고성 콘텐츠로 신뢰도 하락", 1),
        ("② 지역 상권은 외지인 수요를 볼 수 없다", 0),
        ("'서울 사람들이 우리 동네에 관심 있는지' 알 방법 없음", 1),
        ("외지인 수요 데이터는 존재하지만, 상권에 전달할 채널 부재", 1),
        ("③ 지역 소속감을 즐길 콘텐츠가 없다", 0),
        ("서울 vs 부산, 전라 vs 경상 — 건강한 발산 공간 부재", 1),
    ])

    # 4. 솔루션
    add_section_slide(prs, 3, "솔루션")
    add_content_slide(prs, "지역 on 네트워크", [
        "강남on, 송파on, 부산on, 광주on, 여수on, 강릉on 등 시·구 단위 커뮤니티 순차 개설",
        "상위 개념: 서울·경기·충청·전라·강원·경남 등 광역(도) 단위 구성",
        "각 on = 맛집·카페·관광·생활 정보·자유 게시판을 갖춘 독립 커뮤니티",
        "동시에 ONtown 브랜드 아래 하나로 연결",
    ], subtitle="3-1. 지역 on 네트워크")

    add_content_slide(prs, "지역 태그 시스템 — 핵심 차별화", [
        "사용자는 홈 지역(가입 지역)을 보유",
        "강남on 유저가 여수on 활동 → 모든 댓글·게시글에 '강남' 태그 표시",
        ("외지인에게: 현지 생생 추천 + 같은 외지인 후기 태그 필터링", 0),
        ("현지인·상권에게: '강남 사람들이 우리 지역을 보고 있다' 수요 신호", 0),
        ("플랫폼에게: 지역 간 이동·관심 데이터 — 독보적 자산 축적", 0),
    ], subtitle="3-2. 지역 태그 시스템")

    add_two_column_slide(
        prs,
        "클랜 시스템 & 전국 이벤트",
        "🎮 지역 대항 게임 리그",
        [
            "테트리스 등 진입장벽 낮은 게임",
            "지역 예선 → 도(광역) 대회 → 전국 대회",
            "'부산 vs 광주' 구도 = 강력한 관전 콘텐츠",
            "지역 유저 응원 참여 유도",
        ],
        "⭐ 전국 얼짱 대회",
        [
            "지역 예선 → 전국 우승자 선발",
            "제휴 기획사 프로필 제공 → 데뷔 기회",
            "'우리 지역 대표를 스타로' 서사",
            "참여·투표·바이럴 동시 성장",
        ],
    )

    add_content_slide(prs, "이벤트 확장 가능성", [
        "게임·얼짱에서 시작 → 무한 확장",
        "지역 음식 대결 · 사진 공모전 · 지역 밴드 경연",
        "각 이벤트 = 클랜 정체성 강화 + 플랫폼 전체를 축제의 장으로",
    ])

    # 5. 비즈니스 모델
    add_section_slide(prs, 4, "비즈니스 모델")
    add_content_slide(prs, "수익 구조", [
        ("지역 광고 및 제휴 (핵심 수익원)", 0),
        ("맛집·카페·관광 카테고리 지역 소상공인 광고", 1),
        ("'외지인 타겟 광고' — 기존 플랫폼 미제공 상품", 1),
        ("예: 여수 카페 → '강남·송파 유저에게만 노출' 광고", 1),
        ("이벤트 스폰서십", 0),
        ("게임사·음료·식품·지자체 스폰서십", 1),
        ("얼짱 대회 → 뷰티·패션·기획사 제휴", 1),
        ("지자체 B2G + 커머스 연계 (중장기)", 0),
        ("관광 데이터 리포트 · 축제 홍보 · 특산물 공동구매", 1),
    ])

    # 6. 시장 분석
    add_section_slide(prs, 5, "시장 및 경쟁 분석")
    add_two_column_slide(
        prs,
        "시장 기회 & 경쟁 환경",
        "📈 시장 기회",
        [
            "당근마켓 — 하이퍼로컬 잠재력 증명",
            "국내 여행 시장 — 연간 수십조 원",
            "'지역 커뮤니티 × 지역 간 이동' 교차점 = 공백",
        ],
        "⚔️ 경쟁사 한계",
        [
            "당근마켓: 타 지역 탐색·교류 부재",
            "맘카페: 폐쇄적, 외지인 진입 불가",
            "네이버: 광고성 콘텐츠, 커뮤니티성 부재",
        ],
    )
    add_content_slide(prs, "ONtown의 차별화", [
        "'외지인 환영'을 구조화한 최초의 지역 커뮤니티",
        "클랜 기반 이벤트 = 자체 콘텐츠 엔진",
        "지역 태그 = 독보적 데이터 자산 (지역 간 관심·이동)",
        "지역 상권 ↔ 외지인 수요를 연결하는 유일한 채널",
    ])

    # 7. 로드맵
    add_section_slide(prs, 6, "성장 전략")
    add_roadmap_slide(prs)

    # 8. 투자
    add_section_slide(prs, 7, "투자 유치")
    add_funds_slide(prs)

    # 9. 비전
    add_vision_slide(prs)

    # 10. 마무리
    add_title_slide(
        prs,
        "감사합니다",
        "ONtown — 대한민국의 모든 지역을 연결하다",
        "Contact: [이메일 / 연락처를 입력하세요]",
    )

    return prs


if __name__ == "__main__":
    import os
    out_dir = os.path.join(os.path.dirname(__file__), "..", "docs")
    os.makedirs(out_dir, exist_ok=True)
    out_path = os.path.join(out_dir, "ONtown_Business_Plan.pptx")
    prs = build_presentation()
    prs.save(out_path)
    print(f"Saved: {os.path.abspath(out_path)}")
    print(f"Slides: {len(prs.slides)}")
