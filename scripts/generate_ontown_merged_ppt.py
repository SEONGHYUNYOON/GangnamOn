# -*- coding: utf-8 -*-
"""ONtown 통합 IR 사업계획서 PPT — 0711 디자인 + IR 콘텐츠"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ── 0711 브랜드 팔레트 ──
NAVY = RGBColor(0x23, 0x2D, 0x63)
NAVY_MID = RGBColor(0x2F, 0x3C, 0x7E)
CORAL = RGBColor(0xF9, 0x61, 0x67)
YELLOW = RGBColor(0xF9, 0xE7, 0x95)
LIGHT = RGBColor(0xF4, 0xF5, 0xFA)
CREAM = RGBColor(0xFA, 0xF7, 0xF1)
DARK = RGBColor(0x0C, 0x12, 0x20)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRAY = RGBColor(0x64, 0x74, 0x8B)
GRAY_LIGHT = RGBColor(0x94, 0xA3, 0xB8)
SKY = RGBColor(0x0E, 0xA5, 0xE9)
GREEN = RGBColor(0x16, 0xA3, 0x4A)
PURPLE = RGBColor(0x7C, 0x3A, 0xED)
GOLD = RGBColor(0xB5, 0x84, 0x18)
BAR_COLORS = [NAVY_MID, CORAL, GOLD, SKY, GREEN, PURPLE]

FONT = "Malgun Gothic"
W = Inches(13.333)
H = Inches(7.5)


def rgb(c):
    return c


def set_bg(slide, color):
    f = slide.background.fill
    f.solid()
    f.fore_color.rgb = color


def rect(slide, l, t, w, h, fill, line=None, radius=None):
    if radius:
        sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
    else:
        sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    if line:
        sh.line.color.rgb = line
        sh.line.width = Pt(1)
    else:
        sh.line.fill.background()
    return sh


def txt(slide, l, t, w, h, text, size=14, bold=False, color=DARK, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    box = slide.shapes.add_textbox(l, t, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.text = text
    p.font.name = FONT
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = align
    p.line_spacing = 1.25
    return box


def multiline(slide, l, t, w, h, lines, size=13, color=DARK, bullet=False):
    box = slide.shapes.add_textbox(l, t, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = (f"• {line}" if bullet and not line.startswith("•") else line)
        p.font.name = FONT
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.space_after = Pt(6)
        p.line_spacing = 1.3
    return box


def slide_header(slide, section, title, subtitle=None, num=None):
    txt(slide, Inches(0.7), Inches(0.38), Inches(10), Inches(0.3), section, 11, True, CORAL)
    txt(slide, Inches(0.7), Inches(0.68), Inches(11.5), Inches(0.7), title, 26, True, NAVY)
    if subtitle:
        txt(slide, Inches(0.7), Inches(1.35), Inches(11.5), Inches(0.45), subtitle, 12, False, GRAY)
    if num:
        txt(slide, Inches(12.1), Inches(0.35), Inches(0.9), Inches(0.5), f"{num:02d}", 22, True, RGBColor(0xE2, 0xE8, 0xF0), PP_ALIGN.RIGHT)
    rect(slide, Inches(0.7), Inches(0.32), Inches(0.55), Inches(0.04), CORAL)


def footer(slide, n):
    txt(slide, Inches(0.7), Inches(7.05), Inches(6), Inches(0.3), "ONtown Investor Deck", 8, False, GRAY_LIGHT)
    txt(slide, Inches(11.5), Inches(7.05), Inches(1.2), Inches(0.3), f"{n:02d}", 8, False, GRAY_LIGHT, PP_ALIGN.RIGHT)


def pill(slide, l, t, w, h, label, bg=NAVY_MID, fg=WHITE, size=10):
    rect(slide, l, t, w, h, bg, radius=True)
    txt(slide, l, t, w, h, label, size, True, fg, PP_ALIGN.CENTER, MSO_ANCHOR.MIDDLE)


def card(slide, l, t, w, h, title, body, badge=None, icon_color=NAVY_MID):
    rect(slide, l, t, w, h, LIGHT, radius=True)
    rect(slide, l + Inches(0.35), t + Inches(0.35), Inches(0.55), Inches(0.55), icon_color, radius=True)
    txt(slide, l + Inches(0.35), t + Inches(0.85), w - Inches(0.5), Inches(0.45), title, 14, True, NAVY)
    if badge:
        pill(slide, l + w - Inches(1.5), t + Inches(0.35), Inches(1.2), Inches(0.32), badge, CORAL, WHITE, 8)
    txt(slide, l + Inches(0.35), t + Inches(1.35), w - Inches(0.5), h - Inches(1.5), body, 11, False, GRAY)


# ── SLIDES ──

def slide_cover(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s, WHITE)
    # decorative circles (0711 style)
    for l, t, w, h, c in [
        (Inches(9.2), Inches(-1.4), Inches(5.6), Inches(5.6), NAVY),
        (Inches(11.6), Inches(4.9), Inches(3.4), Inches(3.4), CORAL),
        (Inches(8.6), Inches(5.6), Inches(1.1), Inches(1.1), YELLOW),
        (Inches(-1.2), Inches(5.4), Inches(3.2), Inches(3.2), NAVY),
    ]:
        sh = s.shapes.add_shape(MSO_SHAPE.OVAL, l, t, w, h)
        sh.fill.solid()
        sh.fill.fore_color.rgb = c
        sh.line.fill.background()

    txt(s, Inches(0.9), Inches(0.55), Inches(8), Inches(0.35), "LOCAL COMMUNITY NETWORK PLATFORM", 10, True, CORAL)
    txt(s, Inches(0.9), Inches(1.0), Inches(9), Inches(1.2), "ONtown", 52, True, NAVY)
    txt(s, Inches(0.95), Inches(2.15), Inches(9), Inches(0.5), "온타운  ·  하이퍼로컬 커뮤니티 플랫폼", 18, False, NAVY_MID)
    txt(s, Inches(0.95), Inches(2.85), Inches(10), Inches(0.55), "전국 지역 ON을 하나의 세계관으로 연결하는 로컬 커뮤니티 플랫폼", 13, False, GRAY)
    txt(s, Inches(0.95), Inches(3.55), Inches(10.5), Inches(0.8),
        "강남ON에서 시작해 송파ON, 부산ON, 광주ON, 여수ON, 강릉ON으로 확장합니다.\n각 지역은 하나의 클랜이 되고, 사용자는 지역 태그와 미니홈피로 정체성을 쌓습니다.",
        11, False, GRAY)

    regions = ["강남ON", "송파ON", "부산ON", "광주ON", "여수ON", "강릉ON"]
    x = Inches(0.95)
    for r in regions:
        pill(s, x, Inches(4.55), Inches(1.32), Inches(0.48), r)
        x += Inches(1.45)

    rect(s, Inches(0.95), Inches(5.35), Inches(5.5), Inches(0.55), CREAM, radius=True)
    txt(s, Inches(1.1), Inches(5.42), Inches(5.2), Inches(0.4),
        "Seed Investment Proposal · 2026  |  요청 투자금 12억 원  |  목표 손익분기 12개월", 10, True, NAVY_MID, PP_ALIGN.LEFT, MSO_ANCHOR.MIDDLE)
    txt(s, Inches(0.95), Inches(6.55), Inches(8), Inches(0.35), "투자 유치 사업 계획서", 9, False, GRAY_LIGHT)
    footer(s, 1)


def slide_exec(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s, WHITE)
    slide_header(s, "EXECUTIVE SUMMARY", "지역과 지역을 연결하는 전국 커뮤니티 네트워크",
                 "지역 커뮤니티, 여행 정보, 미니홈피, 랭킹 대회를 하나의 계정과 지역 태그로 연결합니다.", 2)

    items = [
        ("01", "지역 = 하나의 클랜", "강남ON, 여수ON, 부산ON처럼 생활권별 커뮤니티를 독립 운영하되 ONtown 브랜드 아래 연결"),
        ("02", "지역 태그", "강남 유저가 여수ON에 댓글을 달면 강남 태그가 표시 — 외지인 수요 데이터가 축적"),
        ("03", "미니홈피", "BGM, 방명록, 일상 피드, 일촌 — 사용자 정체성과 관계가 쌓이는 개인 공간"),
        ("04", "클랜 경쟁", "테트리스·얼짱 대회로 지역 대표와 전국 리그 구성 — 참여·투표·바이럴 동시 유발"),
    ]
    for i, (num, title, desc) in enumerate(items):
        col, row = i % 2, i // 2
        x = Inches(0.7 + col * 6.2)
        y = Inches(1.95 + row * 2.15)
        rect(s, x, y, Inches(5.9), Inches(1.95), LIGHT, radius=True)
        rect(s, x + Inches(0.25), y + Inches(0.25), Inches(0.5), Inches(0.5), CORAL, radius=True)
        txt(s, x + Inches(0.25), y + Inches(0.25), Inches(0.5), Inches(0.5), num, 14, True, WHITE, PP_ALIGN.CENTER, MSO_ANCHOR.MIDDLE)
        txt(s, x + Inches(0.9), y + Inches(0.28), Inches(4.8), Inches(0.4), title, 14, True, NAVY)
        txt(s, x + Inches(0.9), y + Inches(0.75), Inches(4.8), Inches(1.0), desc, 11, False, GRAY)

    rect(s, Inches(0.7), Inches(6.25), Inches(12), Inches(0.55), NAVY, radius=True)
    multiline(s, Inches(0.95), Inches(6.32), Inches(11.5), Inches(0.45), [
        "투자 포인트: 강남ON MVP 기반 지역 복제 확장  ·  지역 간 이동·관심 데이터 누적  ·  광고·ON 재화·스폰서십 수익  ·  12개월 내 월 손익분기"
    ], 10, WHITE)
    footer(s, 2)


def slide_problem(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s, WHITE)
    slide_header(s, "PROBLEM", "왜 지금 ONtown인가 — 시장의 4가지 공백", num=3)

    problems = [
        ("고립된 지역 커뮤니티", "기존 서비스는 거주자만 대상. 여행자·출장자·이주 예정자는 살아있는 정보에 접근 불가. 검색·블로그는 광고성 콘텐츠에 오염."),
        ("수요가 보이지 않는 상권", "여수 카페 사장님은 '서울 사람들이 무엇을 궁금해하는지' 알 방법이 없음. 외지인 수요 데이터는 있지만 상권에 전달할 채널 부재."),
        ("정체성·소속감 부재", "기존 로컬 서비스는 거래·게시글 중심이라 사용자가 누구인지 쌓이지 않음. 지역을 대표하고 경쟁하는 구조가 약함."),
        ("정보 분산", "맛집은 블로그, 관광은 검색, 모임은 오픈채팅, 거래는 당근에 흩어져 있어 '지역 간 이동' 맥락이 사라짐."),
    ]
    for i, (title, body) in enumerate(problems):
        col, row = i % 2, i // 2
        x = Inches(0.7 + col * 6.2)
        y = Inches(1.75 + row * 2.35)
        card(s, x, y, Inches(5.9), Inches(2.15), title, body)

    rect(s, Inches(0.7), Inches(6.45), Inches(12), Inches(0.42), CREAM, radius=True)
    txt(s, Inches(0.95), Inches(6.5), Inches(11.5), Inches(0.35),
        "ONtown은 지역 내부 커뮤니티와 지역 간 이동 데이터를 동시에 잡습니다.", 11, True, NAVY_MID, PP_ALIGN.CENTER, MSO_ANCHOR.MIDDLE)
    footer(s, 3)


def slide_network(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s, WHITE)
    slide_header(s, "SOLUTION 01", "전국 지역 on 네트워크 — 도시에서 전국으로",
                 "시·구 단위 '지역 on' → 광역(도) 커뮤니티 → ONtown 전국 네트워크 3단 구조", 4)

    # top: ONtown hub
    rect(s, Inches(4.9), Inches(1.85), Inches(3.5), Inches(0.95), NAVY, radius=True)
    txt(s, Inches(4.9), Inches(1.95), Inches(3.5), Inches(0.45), "ONtown", 16, True, WHITE, PP_ALIGN.CENTER)
    txt(s, Inches(4.9), Inches(2.35), Inches(3.5), Inches(0.35), "전국 네트워크 · 전국 대회", 10, False, YELLOW, PP_ALIGN.CENTER)

    # mid: regional
    regions = ["서울", "경기", "충청", "전라", "강원", "경남"]
    x = Inches(2.35)
    for r in regions:
        pill(s, x, Inches(3.35), Inches(1.36), Inches(0.55), r, CORAL)
        x += Inches(1.45)
    txt(s, Inches(0.95), Inches(4.05), Inches(11.5), Inches(0.3), "광역 on (도 단위)  —  도 대회 개최 · 광역 교류", 10, False, GRAY, PP_ALIGN.CENTER)

    # bottom: city on
    cities = ["강남ON", "송파ON", "부산ON", "광주ON", "여수ON", "강릉ON", "+ 전국 확장"]
    x = Inches(0.95)
    for i, c in enumerate(cities):
        bg = WHITE if i == 6 else NAVY_MID
        fg = NAVY_MID if i == 6 else WHITE
        pill(s, x, Inches(4.85), Inches(1.58), Inches(0.55), c, bg, fg, 9)
        x += Inches(1.62)
    txt(s, Inches(0.95), Inches(5.55), Inches(11.5), Inches(0.3),
        "지역 on (시·구 단위)  —  맛집 · 카페 · 관광 · 생활 정보 · 모임 · 중고거래", 10, False, GRAY, PP_ALIGN.CENTER)

    # role cards
    roles = [("강남ON", "소속/생활"), ("여수ON", "출장/여행"), ("부산ON", "관광/상권"), ("전국 ONtown", "리그/데이터")]
    for i, (name, role) in enumerate(roles):
        x = Inches(0.7 + i * 3.15)
        rect(s, x, Inches(6.05), Inches(2.9), Inches(0.75), LIGHT, radius=True)
        txt(s, x + Inches(0.15), Inches(6.1), Inches(2.6), Inches(0.3), name, 11, True, NAVY)
        txt(s, x + Inches(0.15), Inches(6.4), Inches(2.6), Inches(0.3), role, 9, False, CORAL)
    footer(s, 4)


def slide_tag_system(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s, WHITE)
    slide_header(s, "SOLUTION 02", "지역 태그 시스템 — ONtown만의 핵심 무기", num=5)

    # mock UI left
    rect(s, Inches(0.7), Inches(1.75), Inches(5.8), Inches(4.85), LIGHT, radius=True)
    txt(s, Inches(0.95), Inches(1.9), Inches(5.3), Inches(0.35), "여수ON  ·  맛집 게시판", 11, True, NAVY)

    # comment 1
    rect(s, Inches(0.95), Inches(2.35), Inches(5.3), Inches(1.15), WHITE, radius=True)
    rect(s, Inches(1.05), Inches(2.5), Inches(0.45), Inches(0.45), NAVY_MID, radius=True)
    txt(s, Inches(1.05), Inches(2.5), Inches(0.45), Inches(0.45), "김", 10, True, WHITE, PP_ALIGN.CENTER, MSO_ANCHOR.MIDDLE)
    txt(s, Inches(1.6), Inches(2.48), Inches(1.2), Inches(0.25), "김O현", 9, True, NAVY)
    pill(s, Inches(2.5), Inches(2.45), Inches(0.65), Inches(0.28), "강남", CORAL, WHITE, 8)
    txt(s, Inches(1.6), Inches(2.78), Inches(4.5), Inches(0.6),
        "내일 여수 출장 갑니다! 밤바다 근처 분위기 좋은 카페 추천 부탁드려요.", 10, False, GRAY)

    # comment 2
    rect(s, Inches(0.95), Inches(3.65), Inches(5.3), Inches(1.15), WHITE, radius=True)
    rect(s, Inches(1.05), Inches(3.8), Inches(0.45), Inches(0.45), NAVY_MID, radius=True)
    txt(s, Inches(1.05), Inches(3.8), Inches(0.45), Inches(0.45), "박", 10, True, WHITE, PP_ALIGN.CENTER, MSO_ANCHOR.MIDDLE)
    txt(s, Inches(1.6), Inches(3.78), Inches(1.2), Inches(0.25), "박O진", 9, True, NAVY)
    pill(s, Inches(2.5), Inches(3.75), Inches(0.65), Inches(0.28), "여수", NAVY_MID, WHITE, 8)
    txt(s, Inches(1.6), Inches(4.08), Inches(4.5), Inches(0.6),
        "현지인 픽으로 낭만포차거리 옆 OO카페 추천해요. 강남 분들이 요즘 많이 오세요!", 10, False, GRAY)

    txt(s, Inches(0.95), Inches(5.0), Inches(5.3), Inches(0.35),
        "타 지역 ON에서 활동하면 홈 지역 태그가 자동 표시됩니다.", 9, True, NAVY_MID)

    # value props right
    props = [
        ("외지인에게", "출장·여행 전 현지 추천 + 같은 외지인 후기 태그 필터링"),
        ("현지 상권에게", "'지금 강남이 우리 지역을 보고 있다' — 수요 신호와 선호 데이터"),
        ("플랫폼에게", "지역 간 이동·관심 데이터 — 어떤 경쟁사도 갖지 못한 독보적 자산"),
    ]
    for i, (title, body) in enumerate(props):
        y = Inches(1.85 + i * 1.55)
        rect(s, Inches(6.85), y, Inches(5.75), Inches(1.35), CREAM if i == 0 else LIGHT, radius=True)
        rect(s, Inches(7.05), y + Inches(0.2), Inches(0.12), Inches(0.12), CORAL)
        txt(s, Inches(7.25), y + Inches(0.15), Inches(5.1), Inches(0.35), title, 13, True, NAVY)
        txt(s, Inches(7.25), y + Inches(0.55), Inches(5.1), Inches(0.7), body, 11, False, GRAY)

    rect(s, Inches(0.7), Inches(6.55), Inches(12), Inches(0.35), NAVY, radius=True)
    txt(s, Inches(0.95), Inches(6.58), Inches(11.5), Inches(0.3),
        "이 댓글 하나가 '외지인 관심' · '강남 유저 취향' · '여수 상권 반응' 데이터로 누적됩니다.", 10, False, WHITE, PP_ALIGN.CENTER, MSO_ANCHOR.MIDDLE)
    footer(s, 5)


def slide_product(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s, WHITE)
    slide_header(s, "PRODUCT", "핵심 제품 구조", num=6)

    modules = [
        ("지역 ON 홈", "맛집·카페·관광·생활정보·모임·중고거래", NAVY_MID),
        ("지역 태그", "댓글/저장/조회에 소속 지역 연결", CORAL),
        ("미니홈피", "BGM, 방명록, 일상 피드, 일촌, 파도타기", SKY),
        ("ON 재화", "꾸미기, 슬롯, 부스팅, 대회 참가", GOLD),
        ("지역 랭킹", "활동 점수, 동네 등급, 클랜 순위", GREEN),
        ("대회/이벤트", "테트리스, 얼짱, 로컬 챌린지", PURPLE),
    ]
    for i, (title, desc, color) in enumerate(modules):
        col, row = i % 3, i // 3
        x = Inches(0.7 + col * 4.15)
        y = Inches(1.85 + row * 2.35)
        rect(s, x, y, Inches(3.85), Inches(2.05), LIGHT, radius=True)
        rect(s, x + Inches(0.25), y + Inches(0.25), Inches(0.45), Inches(0.45), color, radius=True)
        txt(s, x + Inches(0.25), y + Inches(0.85), Inches(3.35), Inches(0.4), title, 14, True, NAVY)
        txt(s, x + Inches(0.25), y + Inches(1.3), Inches(3.35), Inches(0.65), desc, 11, False, GRAY)

    footer(s, 6)


def slide_usecase(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s, WHITE)
    slide_header(s, "USE CASE", "사용자 시나리오: 강남 사람이 여수에 가는 날", num=7)

    steps = [
        ("1", "강남ON 가입", "사용자는 강남 태그와 미니홈피를 가진다"),
        ("2", "여수ON 방문", "출장 전 맛집·카페·관광 정보를 확인한다"),
        ("3", "댓글/저장", "여수 콘텐츠에 강남 태그가 붙은 반응을 남긴다"),
        ("4", "데이터 축적", "여수ON은 외지인 선호와 현지인 추천을 함께 보여준다"),
    ]
    for i, (num, title, desc) in enumerate(steps):
        x = Inches(0.7 + i * 3.1)
        rect(s, x, Inches(1.85), Inches(2.85), Inches(2.5), LIGHT, radius=True)
        rect(s, x + Inches(0.2), Inches(2.05), Inches(0.55), Inches(0.55), CORAL, radius=True)
        txt(s, x + Inches(0.2), Inches(2.05), Inches(0.55), Inches(0.55), num, 18, True, WHITE, PP_ALIGN.CENTER, MSO_ANCHOR.MIDDLE)
        txt(s, x + Inches(0.2), Inches(2.75), Inches(2.45), Inches(0.4), title, 12, True, NAVY)
        txt(s, x + Inches(0.2), Inches(3.2), Inches(2.45), Inches(1.0), desc, 10, False, GRAY)
        if i < 3:
            txt(s, x + Inches(2.95), Inches(2.85), Inches(0.3), Inches(0.4), "→", 16, True, CORAL)

    rect(s, Inches(0.7), Inches(4.65), Inches(5.8), Inches(2.15), NAVY, radius=True)
    txt(s, Inches(0.95), Inches(4.8), Inches(5.3), Inches(0.35), "여수ON에서 가능한 인사이트", 12, True, YELLOW)
    multiline(s, Inches(0.95), Inches(5.2), Inches(5.3), Inches(1.4), [
        "강남 유저가 저장한 여수 맛집 TOP 10",
        "서울권 출장자가 많이 본 카페",
        "현지인 추천 + 외지인 반응 높은 장소",
        "지역별 방문·저장·댓글 전환 데이터",
    ], 10, WHITE, bullet=True)

    rect(s, Inches(6.85), Inches(4.65), Inches(5.85), Inches(2.15), CREAM, radius=True)
    txt(s, Inches(7.1), Inches(4.95), Inches(5.3), Inches(1.5),
        "검색 결과가 아니라\n\"지역 간 반응 데이터\"가\nONtown의 자산입니다.", 16, True, NAVY, PP_ALIGN.LEFT, MSO_ANCHOR.MIDDLE)
    footer(s, 7)


def slide_events(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s, WHITE)
    slide_header(s, "SOLUTION 03", "클랜 시스템과 전국 이벤트 — 성장 엔진",
                 "각 지역 ON은 하나의 클랜. 지역 대항전이 소속감·참여·바이럴을 동시에 폭발시킵니다.", 8)

    # left: game league
    rect(s, Inches(0.7), Inches(1.85), Inches(5.9), Inches(4.75), LIGHT, radius=True)
    txt(s, Inches(0.95), Inches(2.0), Inches(5.4), Inches(0.4), "지역 대항 게임 리그", 15, True, NAVY)
    txt(s, Inches(0.95), Inches(2.45), Inches(5.4), Inches(0.5),
        "테트리스 등 진입장벽 낮은 게임 · \"부산 대표 vs 광주 대표\" 관전 콘텐츠", 10, False, GRAY)
    for i, step in enumerate(["지역 예선", "도(광역) 대회", "전국 대회"]):
        x = Inches(1.1 + i * 1.75)
        pill(s, x, Inches(3.15), Inches(1.45), Inches(0.48), step, NAVY_MID if i < 2 else CORAL)
        if i < 2:
            txt(s, x + Inches(1.5), Inches(3.22), Inches(0.3), Inches(0.35), "→", 14, True, CORAL)
    rect(s, Inches(0.95), Inches(3.95), Inches(5.4), Inches(1.0), NAVY, radius=True)
    txt(s, Inches(0.95), Inches(4.15), Inches(5.4), Inches(0.65), "우승 클랜의 영광\n= 지역 전체의 축제", 13, True, WHITE, PP_ALIGN.CENTER, MSO_ANCHOR.MIDDLE)

    # right: beauty contest
    rect(s, Inches(6.85), Inches(1.85), Inches(5.85), Inches(4.75), LIGHT, radius=True)
    txt(s, Inches(7.1), Inches(2.0), Inches(5.4), Inches(0.4), "전국 얼짱 대회", 15, True, NAVY)
    txt(s, Inches(7.1), Inches(2.45), Inches(5.4), Inches(0.5),
        "\"우리 지역 대표를 스타로\" — 참여·투표·바이럴을 만드는 서사", 10, False, GRAY)
    for i, step in enumerate(["지역 예선·투표", "전국 우승자", "기획사 프로필"]):
        x = Inches(7.25 + i * 1.75)
        pill(s, x, Inches(3.15), Inches(1.45), Inches(0.48), step, NAVY_MID if i < 2 else CORAL)
        if i < 2:
            txt(s, x + Inches(1.5), Inches(3.22), Inches(0.3), Inches(0.35), "→", 14, True, CORAL)
    rect(s, Inches(7.1), Inches(3.95), Inches(5.4), Inches(1.0), CORAL, radius=True)
    txt(s, Inches(7.1), Inches(4.15), Inches(5.4), Inches(0.65), "연예계 데뷔 기회\n= 강력한 참여 동기", 13, True, WHITE, PP_ALIGN.CENTER, MSO_ANCHOR.MIDDLE)

    txt(s, Inches(0.7), Inches(6.55), Inches(12), Inches(0.35),
        "게임·얼짱에서 시작 → 지역 음식 대결 · 사진 공모전 · 지역 밴드 경연으로 무한 확장", 10, False, GRAY, PP_ALIGN.CENTER)
    footer(s, 8)


def slide_growth_loop(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s, WHITE)
    slide_header(s, "GROWTH LOOP", "성장 루프: 콘텐츠 → 정체성 → 경쟁 → 수익", num=9)

    loop = ["콘텐츠", "정체성", "관계", "경쟁", "수익"]
    subs = ["지역 맛집/카페/관광", "미니홈피·지역 태그", "방명록·일촌·대화", "랭킹·대회·리그", "광고·ON재화·스폰서"]
    cx, cy = Inches(6.65), Inches(4.2)
    rect(s, cx - Inches(1.1), cy - Inches(0.55), Inches(2.2), Inches(1.1), NAVY, radius=True)
    txt(s, cx - Inches(1.1), cy - Inches(0.4), Inches(2.2), Inches(0.35), "ONtown", 14, True, WHITE, PP_ALIGN.CENTER)
    txt(s, cx - Inches(1.1), cy - Inches(0.05), Inches(2.2), Inches(0.3), "지역 세계관", 9, False, YELLOW, PP_ALIGN.CENTER)

    positions = [
        (Inches(3.2), Inches(2.2)), (Inches(9.5), Inches(2.2)), (Inches(10.5), Inches(4.5)),
        (Inches(7.5), Inches(6.0)), (Inches(3.5), Inches(6.0)),
    ]
    for i, ((x, y), title, sub) in enumerate(zip(positions, loop, subs)):
        rect(s, x, y, Inches(2.4), Inches(1.05), LIGHT if i % 2 == 0 else CREAM, radius=True)
        rect(s, x + Inches(0.15), y + Inches(0.15), Inches(0.35), Inches(0.35), CORAL, radius=True)
        txt(s, x + Inches(0.6), y + Inches(0.12), Inches(1.65), Inches(0.35), title, 13, True, NAVY)
        txt(s, x + Inches(0.6), y + Inches(0.5), Inches(1.65), Inches(0.45), sub, 9, False, GRAY)
    footer(s, 9)


def slide_business(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s, WHITE)
    slide_header(s, "BUSINESS MODEL", "4개의 수익 축 — 지역 데이터가 곧 상품", num=10)

    models = [
        ("지역 광고 · 제휴", "핵심 수익원", "맛집·카페·관광 소상공인 광고. '외지인 타겟 광고'(강남·송파 유저만 노출)는 독점 상품.", CORAL),
        ("이벤트 스폰서십", "대형 이벤트", "게임 리그 — 게임사·식음료·지자체. 얼짱 — 뷰티·패션·기획사 제휴.", NAVY_MID),
        ("지자체 B2G", "데이터+홍보", "외지인 관심 데이터 리포트, 지역 축제 홍보 채널 판매.", SKY),
        ("ON 재화 · 커머스", "중장기 확장", "미니홈피 꾸미기, 상점 입점, 특산물 공동구매, 대회 굿즈, 클랜 멤버십.", GREEN),
    ]
    for i, (title, badge, body, color) in enumerate(models):
        col, row = i % 2, i // 2
        x = Inches(0.7 + col * 6.2)
        y = Inches(1.85 + row * 2.35)
        rect(s, x, y, Inches(5.9), Inches(2.15), LIGHT, radius=True)
        pill(s, x + Inches(0.3), y + Inches(0.3), Inches(1.1), Inches(0.3), badge, color, WHITE, 8)
        txt(s, x + Inches(0.3), y + Inches(0.75), Inches(5.3), Inches(0.4), title, 14, True, NAVY)
        txt(s, x + Inches(0.3), y + Inches(1.2), Inches(5.3), Inches(0.85), body, 11, False, GRAY)
    footer(s, 10)


def slide_market(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s, WHITE)
    slide_header(s, "MARKET & COMPETITION", "비어 있는 교차점 — 지역 커뮤니티 × 지역 간 이동", num=11)

    rect(s, Inches(0.7), Inches(1.85), Inches(5.8), Inches(4.5), LIGHT, radius=True)
    txt(s, Inches(0.95), Inches(2.0), Inches(5.3), Inches(0.35), "시장 기회", 14, True, NAVY)
    multiline(s, Inches(0.95), Inches(2.45), Inches(5.3), Inches(3.5), [
        "당근마켓 — 하이퍼로컬 잠재력 증명",
        "국내 여행 시장 — 연간 수십조 원 규모",
        "'지역 커뮤니티 × 지역 간 이동' 교차점 = 공백",
        "ONtown = 외지인 환영을 구조화한 최초의 지역 커뮤니티",
    ], 12, DARK, bullet=True)

    competitors = [
        ("당근마켓", "내 동네 중심, 타 지역 탐색·교류 부재"),
        ("맘카페·지역 카페", "폐쇄적, 외지인 진입 불가, UX 파편화"),
        ("네이버·블로그", "광고성 콘텐츠, 커뮤니티성 부재"),
    ]
    for i, (name, desc) in enumerate(competitors):
        y = Inches(1.85 + i * 1.55)
        rect(s, Inches(6.85), y, Inches(5.85), Inches(1.35), CREAM if i == 0 else LIGHT, radius=True)
        txt(s, Inches(7.1), y + Inches(0.15), Inches(5.3), Inches(0.35), name, 13, True, CORAL)
        txt(s, Inches(7.1), y + Inches(0.55), Inches(5.3), Inches(0.7), desc, 11, False, GRAY)

    rect(s, Inches(0.7), Inches(6.45), Inches(12), Inches(0.42), NAVY, radius=True)
    txt(s, Inches(0.95), Inches(6.5), Inches(11.5), Inches(0.35),
        "클랜 이벤트라는 자체 콘텐츠 엔진 + 지역 태그 데이터 = 명확한 차별화", 11, False, WHITE, PP_ALIGN.CENTER, MSO_ANCHOR.MIDDLE)
    footer(s, 11)


def slide_roadmap(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s, WHITE)
    slide_header(s, "ROADMAP", "거점에서 전국으로 — 3단계 성장 전략", num=12)

    phases = [
        ("1", "거점 확보", "0 – 12개월", CORAL, [
            "강남ON 플래그십 런칭",
            "여수·강릉·부산 등 관광 거점 개설",
            "'출장·여행 전 필수 앱' 포지셔닝",
            "0~3M: 강남ON 고도화",
            "4~6M: 서울권 확장 (송파·잠실·성수)",
        ]),
        ("2", "클랜 문화", "12 – 24개월", NAVY_MID, [
            "첫 지역 대항 테트리스 리그",
            "얼짱 대회 지역 예선",
            "광역 ON 오픈 · 지역 태그 고도화",
            "7~9M: 관광 도시 확장",
            "지역 광고 상품 정식 출시",
        ]),
        ("3", "전국 플랫폼", "24개월 ~", NAVY, [
            "전국 대회 정례화",
            "브랜드·기획사 제휴 확대",
            "지자체 B2G 사업 본격화",
            "10~12M: 리그화·스폰서 패키지",
            "커머스 연계 수익 다각화",
        ]),
    ]
    for i, (num, title, period, color, items) in enumerate(phases):
        x = Inches(0.7 + i * 4.15)
        rect(s, x, Inches(1.85), Inches(3.85), Inches(4.85), LIGHT, radius=True)
        rect(s, x, Inches(1.85), Inches(3.85), Inches(0.85), color)
        txt(s, x, Inches(1.95), Inches(0.8), Inches(0.65), num, 28, True, WHITE, PP_ALIGN.CENTER, MSO_ANCHOR.MIDDLE)
        txt(s, x + Inches(0.85), Inches(2.0), Inches(2.8), Inches(0.35), title, 14, True, WHITE)
        txt(s, x + Inches(0.85), Inches(2.38), Inches(2.8), Inches(0.3), period, 10, False, YELLOW)
        multiline(s, x + Inches(0.25), Inches(2.85), Inches(3.35), Inches(3.6), items, 10, GRAY, bullet=True)
    footer(s, 12)


def slide_kpi(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s, WHITE)
    slide_header(s, "KPI & REVENUE", "12개월 핵심 목표 — 12개월 차 월 손익분기", num=13)

    kpis = [
        ("150K", "가입자", "강남·서울권·관광도시"),
        ("60K", "MAU", "콘텐츠/미니홈피 재방문"),
        ("8개", "지역 ON", "강남·송파·여수·강릉·부산 등"),
        ("1.25억", "M12 월매출", "BEP 진입 목표"),
    ]
    for i, (num, label, sub) in enumerate(kpis):
        x = Inches(0.7 + i * 3.15)
        rect(s, x, Inches(1.85), Inches(2.9), Inches(1.65), LIGHT, radius=True)
        txt(s, x, Inches(2.0), Inches(2.9), Inches(0.55), num, 26, True, CORAL, PP_ALIGN.CENTER)
        txt(s, x, Inches(2.6), Inches(2.9), Inches(0.35), label, 12, True, NAVY, PP_ALIGN.CENTER)
        txt(s, x, Inches(2.95), Inches(2.9), Inches(0.35), sub, 9, False, GRAY, PP_ALIGN.CENTER)

    rect(s, Inches(0.7), Inches(3.75), Inches(12), Inches(0.35), CREAM, radius=True)
    txt(s, Inches(0.95), Inches(3.78), Inches(11.5), Inches(0.3),
        "M12 매출 구성: 지역 광고 45M · 상점 입점 28M · ON 재화 24M · 스폰서 18M · 데이터 10M  (단위: 백만 원)", 10, True, NAVY_MID, PP_ALIGN.CENTER, MSO_ANCHOR.MIDDLE)

    rect(s, Inches(0.7), Inches(4.35), Inches(5.8), Inches(2.35), NAVY, radius=True)
    txt(s, Inches(0.95), Inches(4.5), Inches(5.3), Inches(0.35), "12개월 차 재무 목표", 12, True, YELLOW)
    multiline(s, Inches(0.95), Inches(4.95), Inches(5.3), Inches(1.5), [
        "월 매출 1.25억 원",
        "월 비용 1.16억 원",
        "→ 월 손익분기 진입",
        "초기 3개월 제품 완성, 4개월 차부터 매출 발생",
    ], 11, WHITE, bullet=True)

    rect(s, Inches(6.85), Inches(4.35), Inches(5.85), Inches(2.35), LIGHT, radius=True)
    txt(s, Inches(7.1), Inches(4.5), Inches(5.3), Inches(0.35), "추가 KPI", 12, True, NAVY)
    multiline(s, Inches(7.1), Inches(4.95), Inches(5.3), Inches(1.5), [
        "상점 입점 300곳",
        "대회 개최 12회 (지역 예선 + 온라인)",
        "지역 ON 8개 오픈",
    ], 11, GRAY, bullet=True)
    footer(s, 13)


def slide_funding(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s, WHITE)
    slide_header(s, "USE OF FUNDS", "투자금 사용 계획 — 12억 원 / 12개월 런웨이", num=14)

    funds_pct = [
        ("48%", "인건비", "581백만"),
        ("15%", "마케팅/콘텐츠", "180백만"),
        ("10%", "대회/이벤트", "120백만"),
        ("10%", "예비비", "125백만"),
        ("7%", "인프라/API", "84백만"),
        ("6%", "강남 사무실", "74백만"),
        ("3%", "법무/회계", "36백만"),
    ]
    y = Inches(1.8)
    for i, (pct, title, desc) in enumerate(funds_pct):
        pct_val = int(pct.replace("%", ""))
        bar_w = Inches(max(pct_val / 100 * 6.8, 0.45))
        rect(s, Inches(0.7), y, Inches(7.2), Inches(0.52), LIGHT, radius=True)
        rect(s, Inches(0.7), y, bar_w, Inches(0.52), BAR_COLORS[i % len(BAR_COLORS)])
        txt(s, Inches(0.82), y + Inches(0.08), Inches(0.75), Inches(0.35), pct, 12, True, WHITE)
        txt(s, Inches(1.55), y + Inches(0.08), Inches(5.8), Inches(0.35), f"{title}  {desc}", 10, False, WHITE if pct_val > 15 else NAVY)
        y += Inches(0.58)

    rect(s, Inches(8.85), Inches(1.85), Inches(3.85), Inches(5.0), NAVY, radius=True)
    txt(s, Inches(9.1), Inches(2.05), Inches(3.35), Inches(0.4), "총 필요 자금", 12, True, YELLOW)
    txt(s, Inches(9.1), Inches(2.55), Inches(3.35), Inches(0.7), "12억 원", 36, True, WHITE, PP_ALIGN.CENTER)
    txt(s, Inches(9.1), Inches(3.35), Inches(3.35), Inches(0.35), "≈ 1,200백만 원", 11, False, GRAY_LIGHT, PP_ALIGN.CENTER)
    multiline(s, Inches(9.1), Inches(3.85), Inches(3.35), Inches(2.5), [
        "플랫폼 개발 40%",
        "지역 활성화 30%",
        "첫 전국 이벤트 20%",
        "운영·제휴 10%",
        "(사업계획서 기준 배분)",
    ], 10, WHITE, bullet=True)
    footer(s, 14)


def slide_team(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s, WHITE)
    slide_header(s, "TEAM & OFFICE", "초기 8인 팀과 강남 거점 운영", num=15)

    team = [
        ("CEO/PM", "48M", "전략·투자·제품"),
        ("CTO", "90M", "아키텍처·보안·데이터"),
        ("Full-stack ×2", "140M", "웹/앱/관리자"),
        ("Backend/Data", "75M", "지역 태그·랭킹"),
        ("Designer", "55M", "브랜드·UX·미니홈피"),
        ("Growth", "55M", "런칭·콘텐츠·광고"),
        ("Community", "42M", "지역 운영·이벤트"),
    ]
    layouts = [
        (0.7, 1.85, 2.9), (3.55, 1.85, 2.9), (6.4, 1.85, 2.9), (9.25, 1.85, 2.9),
        (2.1, 3.35, 3.2), (5.55, 3.35, 3.2), (9.0, 3.35, 3.2),
    ]
    for (role, pay, desc), (x_in, y_in, w_in) in zip(team, layouts):
        x, y, w = Inches(x_in), Inches(y_in), Inches(w_in)
        rect(s, x, y, w, Inches(1.25), LIGHT, radius=True)
        txt(s, x + Inches(0.15), y + Inches(0.12), w - Inches(0.3), Inches(0.28), role, 11, True, NAVY)
        txt(s, x + Inches(0.15), y + Inches(0.4), Inches(0.9), Inches(0.28), pay, 10, True, CORAL)
        txt(s, x + Inches(0.15), y + Inches(0.68), w - Inches(0.3), Inches(0.45), desc, 9, False, GRAY)

    rect(s, Inches(0.7), Inches(4.85), Inches(12), Inches(1.85), CREAM, radius=True)
    txt(s, Inches(0.95), Inches(5.0), Inches(5.5), Inches(0.35), "강남 사무실 운영 가정", 12, True, NAVY)
    multiline(s, Inches(0.95), Inches(5.35), Inches(5.5), Inches(1.1), [
        "6~8인 업무공간 · 월 450만 원 (보증금·회의 비용 포함)",
        "투자자·파트너 미팅 접근성을 고려한 강남 거점",
    ], 10, GRAY, bullet=True)
    txt(s, Inches(6.85), Inches(5.0), Inches(5.5), Inches(1.5),
        "공개 매물 기준 강남 공유오피스 6인실 월 374만 원 사례 확인.\n본 계획은 여유 좌석·회의/운영 비용을 포함해\n월 450만 원으로 보수적 산정.", 10, False, GRAY)
    footer(s, 15)


def slide_vision(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s, NAVY)
    for l, t, w, h, c in [
        (Inches(10.5), Inches(-1.5), Inches(4.5), Inches(4.5), RGBColor(0x2F, 0x3C, 0x7E)),
        (Inches(-1.5), Inches(5.0), Inches(3.5), Inches(3.5), CORAL),
        (Inches(11.5), Inches(5.5), Inches(1.2), Inches(1.2), YELLOW),
    ]:
        sh = s.shapes.add_shape(MSO_SHAPE.OVAL, l, t, w, h)
        sh.fill.solid()
        sh.fill.fore_color.rgb = c
        sh.line.fill.background()

    txt(s, Inches(0.9), Inches(0.8), Inches(5), Inches(0.35), "INVESTMENT ASK", 11, True, CORAL)
    txt(s, Inches(0.9), Inches(1.3), Inches(11), Inches(1.2),
        "12억 원으로 강남ON을 검증하고\nONtown 전국 확장의 첫 해를 만듭니다", 30, True, WHITE)

    multiline(s, Inches(0.9), Inches(3.0), Inches(6.5), Inches(2.5), [
        "핵심 8인 팀 확보 및 12개월 제품 개발/운영",
        "강남 거점 사무실 · 투자자/파트너 미팅 기반",
        "강남ON 활성화 → 서울권·관광도시 ON 확장",
        "지역 태그 데이터 · 미니홈피 · ON 재화 · 대회 수익화",
    ], 12, RGBColor(0xCB, 0xD5, 0xE1), bullet=True)

    rect(s, Inches(7.8), Inches(2.8), Inches(4.8), Inches(2.8), RGBColor(0x1A, 0x22, 0x45), radius=True)
    txt(s, Inches(8.0), Inches(3.0), Inches(4.4), Inches(0.35), "목표", 11, True, YELLOW)
    txt(s, Inches(8.0), Inches(3.45), Inches(4.4), Inches(0.55), "12개월 차 월 손익분기", 18, True, WHITE)
    txt(s, Inches(8.0), Inches(4.1), Inches(4.4), Inches(0.4), "M12 월매출 1.25억 / 월비용 1.16억", 11, False, GRAY_LIGHT)

    txt(s, Inches(0.9), Inches(5.5), Inches(11.5), Inches(0.8),
        "모든 지역이 클랜으로 살아 숨 쉬고, 서로 경쟁하고 교류하며,\n지역의 상권과 인재가 전국 무대로 올라서는 세계.", 14, False, RGBColor(0xCB, 0xD5, 0xE1))

    txt(s, Inches(0.9), Inches(6.35), Inches(11.5), Inches(0.45),
        "\"당신이 어디에 있든, 그 지역은 ON 되어 있다.\"", 18, True, YELLOW, PP_ALIGN.CENTER)
    txt(s, Inches(0.9), Inches(6.85), Inches(11.5), Inches(0.35),
        "ONtown (온타운)  |  투자 문의: contact@ontown.kr", 10, False, GRAY_LIGHT, PP_ALIGN.CENTER)
    footer(s, 16)


def build():
    prs = Presentation()
    prs.slide_width = W
    prs.slide_height = H

    slide_cover(prs)
    slide_exec(prs)
    slide_problem(prs)
    slide_network(prs)
    slide_tag_system(prs)
    slide_product(prs)
    slide_usecase(prs)
    slide_events(prs)
    slide_growth_loop(prs)
    slide_business(prs)
    slide_market(prs)
    slide_roadmap(prs)
    slide_kpi(prs)
    slide_funding(prs)
    slide_team(prs)
    slide_vision(prs)

    return prs


if __name__ == "__main__":
    import os
    out = os.path.join(os.path.dirname(__file__), "..", "docs", "ONtown_IR_Merged.pptx")
    prs = build()
    prs.save(out)
    print(f"Saved: {os.path.abspath(out)} ({len(prs.slides)} slides)")
