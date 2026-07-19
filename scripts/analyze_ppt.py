"""PPT 분석 스크립트 — 구조·디자인·콘텐츠 추출"""
import json
import os
import sys
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.dml.color import RGBColor

def rgb_str(color):
    try:
        if color and color.type is not None:
            rgb = color.rgb
            if rgb:
                return f"#{rgb}"
    except Exception:
        pass
    return None

def analyze_ppt(path):
    prs = Presentation(path)
    data = {
        "path": path,
        "filename": os.path.basename(path),
        "slide_count": len(prs.slides),
        "width_in": prs.slide_width.inches,
        "height_in": prs.slide_height.inches,
        "slides": [],
    }
    for i, slide in enumerate(prs.slides, 1):
        slide_info = {
            "index": i,
            "layout": slide.slide_layout.name if slide.slide_layout else None,
            "shapes": [],
            "all_text": [],
        }
        for shape in slide.shapes:
            s = {
                "type": str(shape.shape_type),
                "name": shape.name,
                "left_in": round(shape.left.inches, 2) if shape.left else None,
                "top_in": round(shape.top.inches, 2) if shape.top else None,
                "width_in": round(shape.width.inches, 2) if shape.width else None,
                "height_in": round(shape.height.inches, 2) if shape.height else None,
            }
            if shape.has_text_frame:
                s["text"] = shape.text_frame.text
                if s["text"].strip():
                    slide_info["all_text"].append(s["text"].strip())
                paras = []
                for para in shape.text_frame.paragraphs:
                    p = {"text": para.text}
                    if para.font.size:
                        p["size_pt"] = para.font.size.pt
                    if para.font.bold:
                        p["bold"] = True
                    try:
                        if para.font.color and para.font.color.rgb:
                            p["color"] = rgb_str(para.font.color)
                    except Exception:
                        pass
                    if para.text.strip():
                        paras.append(p)
                s["paragraphs"] = paras
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                s["is_picture"] = True
            try:
                if hasattr(shape, "fill") and shape.fill.type is not None:
                    s["fill_color"] = rgb_str(shape.fill.fore_color)
            except Exception:
                pass
            slide_info["shapes"].append(s)
        data["slides"].append(slide_info)
    return data


if __name__ == "__main__":
    files = [
        r"d:\웨일 다운로드\ONtown_사업계획서_0711.pptx",
        r"d:\DEV\GangNam_On\docs\ONtown_IR_사업계획서.pptx",
    ]
    out = []
    for fp in files:
        if os.path.exists(fp):
            out.append(analyze_ppt(fp))
        else:
            out.append({"path": fp, "error": "not found"})
    out_path = r"D:\DEV\GangNam_On\docs\ppt_analysis.json"
    with open(out_path, "w", encoding="utf-8") as f:
        json.dump(out, f, ensure_ascii=False, indent=2)
    print(out_path)
